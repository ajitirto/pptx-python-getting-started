from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

# Membuat presentasi baru
prs = Presentation()

# Menambahkan slide dengan layout 'Title Slide' untuk halaman intro
slide_layout = prs.slide_layouts[0]  # Title slide layout
slide = prs.slides.add_slide(slide_layout)

# Menambahkan judul
title = slide.shapes.title
title.text = "Introduction"

# Menambahkan subtitle dengan tanggal saat ini
subtitle = slide.placeholders[1]
subtitle.text = f"Today's date: {datetime.now().strftime('%d %B %Y')}"

# Menambahkan style menarik untuk judul
title_text_frame = title.text_frame
title_text_frame.paragraphs[0].font.size = Pt(48)
title_text_frame.paragraphs[0].font.bold = True
title_text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 70, 130)  # Biru tua

# Menambahkan style menarik untuk subtitle
subtitle_text_frame = subtitle.text_frame
subtitle_text_frame.paragraphs[0].font.size = Pt(24)
subtitle_text_frame.paragraphs[0].font.color.rgb = RGBColor(90, 90, 90)  # Abu-abu
subtitle_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Menyimpan presentasi
prs.save('Introduction_Presentation.pptx')

print("Presentasi berhasil dibuat!")
